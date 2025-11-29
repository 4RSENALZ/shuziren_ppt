from __future__ import annotations

import base64
import json
import os
import re
import time
import uuid
import shutil
import tempfile
import subprocess
from fractions import Fraction
from threading import Event, Thread
from typing import Callable, Dict, List, Tuple, Optional, TYPE_CHECKING, Any

from docx import Document
from pptx import Presentation
from pptx.util import Inches

import cv2

from logger import logger

try:
    import comtypes
    import comtypes.client
except ImportError:  # pragma: no cover - optional dependency
    comtypes = None

if TYPE_CHECKING:
    from basereal import BaseReal
else:  # pragma: no cover - typing helper
    BaseReal = object  # type: ignore[assignment]


class DigitalHumanRenderer:
    """Helper to reuse a single digital human instance for multiple scripts."""

    def __init__(
        self,
        session_factory: Callable[[], int],
        builder: Callable[[int], BaseReal],
        warmup_delay: float = 2.0,
        speak_timeout: Optional[float] = None,
        max_preroll: float = 1.5,
    ) -> None:
        from basereal import BaseReal  # Local import to avoid cycles

        self._session_factory = session_factory
        self._builder = builder
        self._speak_timeout = speak_timeout if speak_timeout and speak_timeout > 0 else None
        self._max_preroll = max(0.0, float(max_preroll))
        self._start_event = Event()
        self._end_event = Event()
        self._quit_event = Event()
        self._ffmpeg_bin = self._locate_ffmpeg_binary()

        self.session_id = self._session_factory()
        self._nerfreal: BaseReal = self._builder(self.session_id)

        original_notify = getattr(self._nerfreal, "notify", None)

        def patched_notify(eventpoint):
            if callable(original_notify):
                original_notify(eventpoint)
            if not isinstance(eventpoint, dict):
                return
            status = eventpoint.get("status")
            if status == "start":
                self._start_event.set()
            elif status == "end":
                self._end_event.set()

        self._nerfreal.notify = patched_notify  # type: ignore[attr-defined]

        self._render_thread = Thread(
            target=self._nerfreal.render,
            args=(self._quit_event,),
            daemon=True,
        )
        self._render_thread.start()

        warmup_deadline = time.time() + max(0.5, warmup_delay)
        while (self._nerfreal.width == 0 or self._nerfreal.height == 0) and time.time() < warmup_deadline:
            time.sleep(0.1)

    def speak(self, text: str, output_path: str) -> str:
        if not text.strip():
            raise ValueError("朗读内容不能为空。")

        wait_deadline = time.time() + 20.0
        while (self._nerfreal.width == 0 or self._nerfreal.height == 0) and time.time() < wait_deadline:
            time.sleep(0.1)
        if self._nerfreal.width == 0 or self._nerfreal.height == 0:
            raise RuntimeError("数字人视频流未就绪，无法开始录制。")

        self._start_event.clear()
        self._end_event.clear()

        record_start = time.time()
        self._nerfreal.start_recording(output_path)
        is_ssml = text.lstrip().lower().startswith("<speak")
        meta = {"source": "ppt"}
        if is_ssml:
            meta["ssml"] = True
        self._nerfreal.put_msg_txt(text, meta)

        if self._speak_timeout is not None and not self._start_event.wait(timeout=self._speak_timeout):
            self._nerfreal.stop_recording()
            raise TimeoutError("数字人朗读未按预期启动，请稍后重试。")

        event_start_time = time.time()

        speak_deadline = None if self._speak_timeout is None else time.time() + self._speak_timeout
        poll_interval = 0.05  # shorter poll for quicker start/stop detection
        speech_start_time: Optional[float] = None
        while True:
            if speak_deadline is not None and time.time() >= speak_deadline:
                self._nerfreal.stop_recording()
                raise TimeoutError("数字人朗读未在限定时间内完成。")
            speaking_now = self._nerfreal.is_speaking()
            if speaking_now and speech_start_time is None:
                speech_start_time = time.time()
            if self._end_event.is_set() and not speaking_now:
                break
            time.sleep(poll_interval)

        # 等待音视频管线彻底排空
        drain_deadline = time.time() + 3.0
        while self._nerfreal.is_speaking() and time.time() < drain_deadline:
            time.sleep(poll_interval)
        time.sleep(0.1)

        self._nerfreal.stop_recording()
        self._nerfreal.flush_talk()

        effective_start = speech_start_time or event_start_time
        wait_estimate = max(0.0, effective_start - record_start)
        detected_wait = self._detect_audio_lead(output_path)
        wait_before_speech = detected_wait if detected_wait is not None else wait_estimate
        self._trim_preroll(output_path, wait_before_speech)
        return output_path

    def close(self) -> None:
        self._quit_event.set()
        if self._render_thread.is_alive():
            self._render_thread.join(timeout=5.0)

    def _locate_ffmpeg_binary(self) -> Optional[str]:
        root_dir = os.path.dirname(os.path.abspath(__file__))
        candidates = [
            os.path.join(root_dir, "ffmpeg", "ffmpeg.exe"),
            os.path.join(root_dir, "ffmpeg", "bin", "ffmpeg.exe"),
            os.path.join(root_dir, "ffmpeg", "ffmpeg"),
            os.path.join(root_dir, "ffmpeg", "bin", "ffmpeg"),
        ]
        for candidate in candidates:
            if os.path.isfile(candidate):
                return candidate
        return shutil.which("ffmpeg") or shutil.which("ffmpeg.exe")

    def _detect_audio_lead(self, video_path: str) -> Optional[float]:
        if not self._ffmpeg_bin or not os.path.isfile(video_path):
            return None

        command = [
            self._ffmpeg_bin,
            "-hide_banner",
            "-nostats",
            "-i",
            video_path,
            "-vn",
            "-af",
            "silencedetect=noise=-40dB:d=0.05",
            "-f",
            "null",
            "-",
        ]
        try:
            result = subprocess.run(
                command,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
                check=False,
            )
        except Exception as exc:
            logger.warning("[ppt-augment] detect audio lead failed: %s", exc)
            return None

        pattern = re.compile(r"silence_end:\s*([0-9]+(?:\.[0-9]+)?)")
        for line in result.stderr.splitlines():
            match = pattern.search(line)
            if match:
                try:
                    return max(0.0, float(match.group(1)))
                except ValueError:
                    continue
        return None

    def _trim_preroll(self, output_path: str, wait_before_speech: float) -> None:
        if self._max_preroll <= 0.0:
            return
        tolerance = 0.005
        if not self._ffmpeg_bin or not os.path.isfile(output_path):
            logger.warning("[ppt-augment] ffmpeg unavailable, skip preroll trimming")
            return

        current_wait = max(0.0, float(wait_before_speech))
        for attempt in range(3):
            if current_wait <= self._max_preroll + tolerance:
                break
            overshoot = 0.01
            trim_offset = max(0.0, current_wait - self._max_preroll + overshoot)
            if trim_offset <= 0.0:
                break

            tmp_fd, tmp_path = tempfile.mkstemp(suffix=".mp4")
            os.close(tmp_fd)
            command = [
                self._ffmpeg_bin,
                "-y",
                "-i",
                output_path,
                "-ss",
                f"{trim_offset:.3f}",
                "-map",
                "0:v:0",
                "-map",
                "0:a?",
                "-c:v",
                "libx264",
                "-preset",
                "veryfast",
                "-crf",
                "18",
                "-pix_fmt",
                "yuv420p",
                "-c:a",
                "aac",
                "-b:a",
                "192k",
                "-ar",
                "48000",
                "-movflags",
                "+faststart",
                "-reset_timestamps",
                "1",
                tmp_path,
            ]
            try:
                subprocess.run(command, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
            except subprocess.CalledProcessError as exc:
                logger.warning("[ppt-augment] trim preroll failed: %s", exc)
                try:
                    os.remove(tmp_path)
                except OSError:
                    pass
                break

            try:
                os.replace(tmp_path, output_path)
            except Exception as exc:
                logger.warning("[ppt-augment] replace trimmed video failed: %s", exc)
                try:
                    os.remove(tmp_path)
                except OSError:
                    pass
                break

            detected_wait = self._detect_audio_lead(output_path)
            if detected_wait is not None:
                current_wait = max(0.0, detected_wait)
            else:
                current_wait = max(0.0, current_wait - trim_offset)



def generate_slide_previews(ppt_bytes: bytes, limit: int = 0) -> Tuple[List[str], Tuple[float, float]]:
    if comtypes is None:
        raise RuntimeError('需要安装 comtypes 并确保系统已安装 Microsoft PowerPoint 才能生成 PPT 预览。')

    workdir = tempfile.mkdtemp(prefix="pptpreview_")
    try:
        ppt_path = os.path.join(workdir, "input.pptx")
        with open(ppt_path, "wb") as f:
            f.write(ppt_bytes)

        ppt_path = os.path.abspath(ppt_path)

        export_dir = os.path.join(workdir, "slides")
        os.makedirs(export_dir, exist_ok=True)

        logger.info("[ppt-preview] start export for %s (limit=%s)", ppt_path, limit)
        comtypes.CoInitialize()
        powerpoint = None
        presentation = None
        try:
            powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
            powerpoint.Visible = 1
            presentation = powerpoint.Presentations.Open(ppt_path, WithWindow=True)
            try:
                slide_width_pts = presentation.PageSetup.SlideWidth
                slide_height_pts = presentation.PageSetup.SlideHeight
                presentation.Export(export_dir, "PNG")
            finally:
                presentation.Close()
        finally:
            if powerpoint is not None:
                try:
                    powerpoint.Quit()
                except Exception:
                    pass
            comtypes.CoUninitialize()

        slide_width_in = float(slide_width_pts) / 72.0
        slide_height_in = float(slide_height_pts) / 72.0

        def _slide_order_key(entry: str) -> tuple[int, str]:
            match = re.search(r"(\d+)(?=\.[^.]+$)", entry)
            if match:
                try:
                    return int(match.group(1)), entry
                except ValueError:
                    pass
            return (10**9, entry)

        images = []
        ordered_files = sorted(os.listdir(export_dir), key=_slide_order_key)
        for idx, filename in enumerate(ordered_files):
            if limit and idx >= limit:
                break
            file_path = os.path.join(export_dir, filename)
            with open(file_path, "rb") as img_file:
                images.append(base64.b64encode(img_file.read()).decode("ascii"))

        if not images:
            raise RuntimeError("未能导出 PPT 预览图，请确认已安装 Microsoft PowerPoint。")

        logger.info(
            "[ppt-preview] done: %d slide images (~%.2f\" x %.2f\")",
            len(images),
            slide_width_in,
            slide_height_in,
        )
        return images, (slide_width_in, slide_height_in)
    finally:
        shutil.rmtree(workdir, ignore_errors=True)


class PPTDigitalHumanAugmenter:
    def __init__(
        self,
        session_factory: Callable[[], int],
        builder: Callable[[int], BaseReal],
        video_position: Tuple[float, float, float, float] | None = None,
    ) -> None:
        self._session_factory = session_factory
        self._builder = builder
        if video_position is None:
            # (left, top, width, height) in inches
            self._video_position = (5.0, 1.0, 4.5, 3.2)
        else:
            self._video_position = video_position
        self._ffmpeg_bin = self._locate_ffmpeg_binary()
        self._course_base_height = 1080
        self._course_static_duration = 3.0
        self._session_store_root = os.path.join(tempfile.gettempdir(), "pptaugment_sessions")
        os.makedirs(self._session_store_root, exist_ok=True)
        self._active_sessions: Dict[str, Dict[str, Any]] = {}
        self._max_sessions = 8

    def _parse_scripts(self, docx_path: str) -> Tuple[Dict[int, str], Dict[str, str]]:
        document = Document(docx_path)
        pattern = re.compile(r"^p(\d+)[\s:：\.-]*", re.IGNORECASE)
        scripts: Dict[int, List[str]] = {}
        current_slide = None
        ssml_prefix: Optional[str] = None
        ssml_suffix: Optional[str] = None
        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
            lowered = text.lower()
            if (
                ssml_prefix is None
                and not scripts
                and current_slide is None
                and lowered.startswith("<speak")
            ):
                ssml_prefix = text
                continue
            if lowered == "</speak>":
                ssml_suffix = text
                continue
            match = pattern.match(text)
            if match:
                slide_index = int(match.group(1))
                content = text[match.end():].strip()
                current_slide = slide_index
                scripts[current_slide] = []
                if content:
                    scripts[current_slide].append(content)
            elif current_slide is not None:
                scripts[current_slide].append(text)
        flattened = {idx: "\n".join(lines).strip() for idx, lines in scripts.items() if lines}
        meta: Dict[str, str] = {}
        if ssml_prefix:
            meta["ssml_prefix"] = ssml_prefix
        if ssml_suffix:
            meta["ssml_suffix"] = ssml_suffix
        return flattened, meta

    def _locate_ffmpeg_binary(self) -> Optional[str]:
        root_dir = os.path.dirname(os.path.abspath(__file__))
        candidates = [
            os.path.join(root_dir, "ffmpeg", "ffmpeg.exe"),
            os.path.join(root_dir, "ffmpeg", "bin", "ffmpeg.exe"),
            os.path.join(root_dir, "ffmpeg", "ffmpeg"),
            os.path.join(root_dir, "ffmpeg", "bin", "ffmpeg"),
        ]
        for candidate in candidates:
            if os.path.isfile(candidate):
                return candidate
        return shutil.which("ffmpeg") or shutil.which("ffmpeg.exe")

    def _ensure_video_aspect_ratio(self, video_path: str, target_ratio: Optional[float]) -> str:
        if not target_ratio or target_ratio <= 0:
            return video_path

        cap = cv2.VideoCapture(video_path)
        if not cap.isOpened():
            return video_path
        try:
            width = cap.get(cv2.CAP_PROP_FRAME_WIDTH) or 0
            height = cap.get(cv2.CAP_PROP_FRAME_HEIGHT) or 0
        finally:
            cap.release()

        if width <= 0 or height <= 0:
            return video_path

        actual_ratio = float(width) / float(height)
        if abs(actual_ratio - target_ratio) <= 0.01:
            return video_path

        if not self._ffmpeg_bin:
            logger.warning("[ppt-augment] ffmpeg not available, skip aspect adjustment")
            return video_path

        fraction = Fraction(target_ratio).limit_denominator(1000)
        aspect_expr = f"{fraction.numerator}:{fraction.denominator}"
        base_dir = os.path.dirname(os.path.abspath(video_path)) or None
        tmp_fd, adjusted_path = tempfile.mkstemp(
            dir=base_dir,
            prefix=f"{os.path.splitext(os.path.basename(video_path))[0]}_dar_",
            suffix=".mp4",
        )
        os.close(tmp_fd)
        command = [
            self._ffmpeg_bin,
            "-y",
            "-i",
            video_path,
            "-vf",
            f"setsar=1,setdar={aspect_expr}",
            "-c:v",
            "libx264",
            "-preset",
            "veryfast",
            "-crf",
            "18",
            "-pix_fmt",
            "yuv420p",
            "-c:a",
            "aac",
            "-b:a",
            "192k",
            "-ar",
            "48000",
            "-movflags",
            "+faststart",
            adjusted_path,
        ]
        try:
            subprocess.run(command, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        except subprocess.CalledProcessError as exc:
            logger.warning("[ppt-augment] adjust video aspect failed: %s", exc)
            if os.path.exists(adjusted_path):
                try:
                    os.remove(adjusted_path)
                except OSError:
                    pass
            return video_path

        try:
            shutil.move(adjusted_path, video_path)
        except Exception as exc:
            logger.warning("[ppt-augment] replace adjusted video failed: %s", exc)
            if os.path.exists(adjusted_path):
                try:
                    os.remove(adjusted_path)
                except OSError:
                    pass
            return video_path

        return video_path

    def _cleanup_session(self, session_id: str) -> None:
        record = self._active_sessions.pop(session_id, None)
        session_dir = None
        if record and isinstance(record, dict):
            session_dir = record.get("path")  # type: ignore[assignment]
        if not session_dir:
            session_dir = os.path.join(self._session_store_root, session_id)
        shutil.rmtree(session_dir, ignore_errors=True)

    def _register_session(self, session_id: str, session_dir: str) -> None:
        self._active_sessions[session_id] = {
            "path": session_dir,
            "created": time.time(),
        }
        while len(self._active_sessions) > self._max_sessions:
            oldest_id = min(
                self._active_sessions.items(),
                key=lambda item: item[1].get("created", 0.0),
            )[0]
            if oldest_id == session_id:
                break
            self._cleanup_session(oldest_id)

    def _save_session_meta(self, session_dir: str, meta: Dict[str, Any]) -> None:
        meta_path = os.path.join(session_dir, "meta.json")
        with open(meta_path, "w", encoding="utf-8") as meta_file:
            json.dump(meta, meta_file, ensure_ascii=False, indent=2)

    def _persist_session_assets(
        self,
        ppt_bytes: bytes,
        generated_videos: Dict[int, str],
        applied_positions: Dict[int, Dict[str, float]],
        slide_dims_in: Tuple[float, float],
        total_slides: int,
        voice_id: Optional[str] = None,
        slide_durations: Optional[Dict[int, float]] = None,
        slide_cues: Optional[Dict[int, float]] = None,
        combined_video_path: Optional[str] = None,
        final_video_path: Optional[str] = None,
        final_video_meta: Optional[Dict[str, Any]] = None,
    ) -> str:
        session_id = uuid.uuid4().hex
        session_dir = os.path.join(self._session_store_root, session_id)
        if os.path.isdir(session_dir):
            shutil.rmtree(session_dir, ignore_errors=True)
        os.makedirs(session_dir, exist_ok=True)

        ppt_filename = "original.pptx"
        with open(os.path.join(session_dir, ppt_filename), "wb") as f:
            f.write(ppt_bytes)

        videos_dir = os.path.join(session_dir, "videos")
        os.makedirs(videos_dir, exist_ok=True)
        video_index: Dict[str, str] = {}
        for slide_idx, src_path in generated_videos.items():
            if not os.path.isfile(src_path):
                continue
            dest_name = f"slide_{slide_idx:03d}.mp4"
            dest_path = os.path.join(videos_dir, dest_name)
            shutil.copy2(src_path, dest_path)
            video_index[str(slide_idx)] = dest_name

        combined_filename = None
        if combined_video_path and os.path.isfile(combined_video_path):
            combined_filename = "combined.mp4"
            combined_target = os.path.join(session_dir, combined_filename)
            shutil.copy2(combined_video_path, combined_target)

        final_filename = None
        if final_video_path and os.path.isfile(final_video_path):
            final_filename = "course_final.mp4"
            shutil.copy2(final_video_path, os.path.join(session_dir, final_filename))

        meta: Dict[str, Any] = {
            "total_slides": total_slides,
            "applied_positions": applied_positions,
            "slide_dims": list(slide_dims_in),
            "videos": video_index,
            "ppt_filename": ppt_filename,
            "created_at": time.time(),
            "static_duration": self._course_static_duration,
            "voice": voice_id,
            "slide_durations": {str(k): float(v) for k, v in (slide_durations or {}).items()},
            "slide_cues": {str(k): float(v) for k, v in (slide_cues or {}).items()},
            "combined_video": combined_filename,
            "final_video": final_filename,
            "course": {
                "ready": False,
                "filename": "course.mp4",
                "duration": None,
                "segments": None,
                "resolution": None,
            },
        }

        if final_filename and final_video_meta:
            course_info = meta.get("course", {})
            course_info.update(
                {
                    "ready": True,
                    "filename": final_filename,
                    "duration": float(final_video_meta.get("duration", 0.0)),
                    "segments": 1,
                    "resolution": list(final_video_meta.get("resolution", (0, 0))),
                    "generated_at": time.time(),
                    "position": final_video_meta.get("position"),
                    "source": "overlay",
                }
            )
            meta["course"] = course_info

        self._save_session_meta(session_dir, meta)
        self._register_session(session_id, session_dir)
        return session_id

    def _load_session_assets(
        self,
        session_id: str,
    ) -> Tuple[bytes, Dict[int, str], Dict[int, Dict[str, float]], Tuple[float, float], int, str, Dict[str, Any]]:
        session_dir = os.path.join(self._session_store_root, session_id)
        if not os.path.isdir(session_dir):
            raise KeyError(f"未找到会话 {session_id}")

        meta_path = os.path.join(session_dir, "meta.json")
        if not os.path.isfile(meta_path):
            raise KeyError(f"会话 {session_id} 的元数据缺失")

        with open(meta_path, "r", encoding="utf-8") as meta_file:
            meta = json.load(meta_file)

        ppt_filename = meta.get("ppt_filename", "original.pptx")
        ppt_path = os.path.join(session_dir, ppt_filename)
        if not os.path.isfile(ppt_path):
            raise KeyError(f"会话 {session_id} 的 PPT 数据缺失")

        with open(ppt_path, "rb") as ppt_file:
            ppt_bytes = ppt_file.read()

        videos_dir = os.path.join(session_dir, "videos")
        video_index_raw = meta.get("videos", {}) or {}
        videos: Dict[int, str] = {}
        for key, filename in video_index_raw.items():
            try:
                slide_idx = int(key)
            except (TypeError, ValueError):
                continue
            candidate_path = os.path.join(videos_dir, filename)
            if os.path.isfile(candidate_path):
                videos[slide_idx] = candidate_path

        applied_positions_raw = meta.get("applied_positions", {}) or {}
        applied_positions: Dict[int, Dict[str, float]] = {}
        for key, position in applied_positions_raw.items():
            try:
                slide_idx = int(key)
            except (TypeError, ValueError):
                continue
            if isinstance(position, dict):
                applied_positions[slide_idx] = {
                    axis: float(position.get(axis, 0.0))
                    for axis in ("x", "y", "width", "height")
                }

        slide_dims_raw = meta.get("slide_dims", [0.0, 0.0])
        slide_dims_in = (
            float(slide_dims_raw[0]) if len(slide_dims_raw) > 0 else 0.0,
            float(slide_dims_raw[1]) if len(slide_dims_raw) > 1 else 0.0,
        )
        total_slides = int(meta.get("total_slides", max(videos.keys(), default=0)))

        self._register_session(session_id, session_dir)

        return ppt_bytes, videos, applied_positions, slide_dims_in, total_slides, session_dir, meta

    def compose_course_from_session(
        self,
        session_id: str,
        course_filename: str | None = None,
        static_duration: Optional[float] = None,
    ) -> Tuple[bytes, Dict[str, object]]:
        ppt_bytes, videos, applied_positions, slide_dims_in, total_slides, session_dir, meta = self._load_session_assets(session_id)

        if not self._ffmpeg_bin:
            raise RuntimeError("未检测到 ffmpeg，无法生成视频课。")

        course_info = meta.get("course") or {}
        desired_filename = course_filename or course_info.get("filename") or "course.mp4"
        stored_course_filename = os.path.basename(desired_filename)
        course_path = os.path.join(session_dir, stored_course_filename)
        if os.path.isfile(course_path):
            with open(course_path, "rb") as course_file:
                course_bytes = course_file.read()
            stored_resolution = course_info.get("resolution")
            if stored_resolution and isinstance(stored_resolution, list):
                cached_resolution = (stored_resolution[0], stored_resolution[1])
            else:
                cached_resolution = tuple(stored_resolution) if stored_resolution else (0, 0)

            return course_bytes, {
                "filename": stored_course_filename,
                "duration": course_info.get("duration"),
                "segments": course_info.get("segments"),
                "resolution": cached_resolution,
            }

        if static_duration is not None and static_duration > 0:
            original_duration = self._course_static_duration
            self._course_static_duration = static_duration
        else:
            original_duration = None

        workdir = tempfile.mkdtemp(prefix=f"course_{session_id}_")
        try:
            final_filename = course_filename or stored_course_filename
            course_bytes, course_meta = self._compose_course_video(
                ppt_bytes,
                videos,
                applied_positions,
                slide_dims_in,
                total_slides,
                workdir,
                final_filename,
            )

            target_path = os.path.join(session_dir, os.path.basename(final_filename))
            with open(target_path, "wb") as course_file:
                course_file.write(course_bytes)

            stored_course_filename = os.path.basename(final_filename)
            course_path = target_path
            resolution = course_meta.get("resolution")
            if resolution and isinstance(resolution, tuple):
                resolution_value: Any = list(resolution)
            else:
                resolution_value = resolution

            course_info.update(
                {
                    "ready": True,
                    "filename": stored_course_filename,
                    "duration": course_meta.get("duration"),
                    "segments": course_meta.get("segments"),
                    "resolution": resolution_value,
                    "generated_at": time.time(),
                }
            )
            meta["course"] = course_info
            self._save_session_meta(session_dir, meta)

            return course_bytes, course_meta
        finally:
            shutil.rmtree(workdir, ignore_errors=True)
            if original_duration is not None:
                self._course_static_duration = original_duration

    def _embed_video(self, slide, video_path: str, position: Optional[Dict[str, float]], slide_dims_in: Tuple[float, float]) -> Dict[str, float]:
        slide_width_in, slide_height_in = slide_dims_in
        target_ratio = None
        if position:
            left_ratio = max(0.0, min(1.0, position.get("x", 0.0)))
            top_ratio = max(0.0, min(1.0, position.get("y", 0.0)))
            width_ratio = position.get("width")
            height_ratio = position.get("height")
            default_width_ratio = self._video_position[2] / slide_width_in if slide_width_in else 0.3
            default_height_ratio = self._video_position[3] / slide_height_in if slide_height_in else 0.3
            if width_ratio is None or width_ratio <= 0:
                width_ratio = default_width_ratio
            if height_ratio is None or height_ratio <= 0:
                height_ratio = default_height_ratio
            width_ratio = max(0.02, min(1.0, width_ratio))
            height_ratio = max(0.02, min(1.0, height_ratio))
            left_in = slide_width_in * min(left_ratio, 1.0 - width_ratio)
            top_in = slide_height_in * min(top_ratio, 1.0 - height_ratio)
            width_in = slide_width_in * width_ratio
            height_in = slide_height_in * height_ratio
        else:
            left_in, top_in, width_in, height_in = self._video_position

        if height_in > 0:
            target_ratio = width_in / height_in

        video_path = self._ensure_video_aspect_ratio(video_path, target_ratio)
        movie = slide.shapes.add_movie(
            video_path,
            Inches(left_in),
            Inches(top_in),
            Inches(width_in),
            Inches(height_in),
            mime_type="video/mp4",
        )
        try:
            movie.lock_aspect_ratio = False
            movie.width = Inches(width_in)
            movie.height = Inches(height_in)
        except Exception:
            pass
        used_position = {
            "x": min(1.0, max(0.0, left_in / slide_width_in)) if slide_width_in else 0.0,
            "y": min(1.0, max(0.0, top_in / slide_height_in)) if slide_height_in else 0.0,
            "width": min(1.0, max(0.0, width_in / slide_width_in)) if slide_width_in else 0.0,
            "height": min(1.0, max(0.0, height_in / slide_height_in)) if slide_height_in else 0.0,
        }
        return used_position

    def _build_course_segment_with_avatar(
        self,
        background_path: str,
        video_path: str,
        position: Dict[str, float],
        target_width: int,
        target_height: int,
        workdir: str,
        slide_idx: int,
    ) -> Tuple[str, float]:
        overlay_x = int(round(position.get("x", 0.0) * target_width))
        overlay_y = int(round(position.get("y", 0.0) * target_height))
        overlay_w = int(round(position.get("width", 0.0) * target_width))
        overlay_h = int(round(position.get("height", 0.0) * target_height))
        overlay_w = max(2, min(target_width, overlay_w))
        overlay_h = max(2, min(target_height, overlay_h))
        if overlay_w % 2 != 0:
            overlay_w += 1
        if overlay_h % 2 != 0:
            overlay_h += 1
        max_x = max(0, target_width - overlay_w)
        max_y = max(0, target_height - overlay_h)
        overlay_x = max(0, min(max_x, overlay_x))
        overlay_y = max(0, min(max_y, overlay_y))

        duration = 0.0
        cap = cv2.VideoCapture(video_path)
        if cap.isOpened():
            try:
                fps = cap.get(cv2.CAP_PROP_FPS) or 0.0
                frame_count = cap.get(cv2.CAP_PROP_FRAME_COUNT) or 0.0
                if fps > 0.0 and frame_count > 0.0:
                    duration = float(frame_count / fps)
            finally:
                cap.release()
        else:
            cap.release()

        segment_path = os.path.join(workdir, f"course_seg_{slide_idx:03d}.mp4")
        filter_complex = (
            f"[0:v]scale={target_width}:{target_height}[bg];"
            f"[1:v]scale={overlay_w}:{overlay_h}:force_original_aspect_ratio=increase,"
            f"crop={overlay_w}:{overlay_h},setsar=1[avatar];"
            f"[bg][avatar]overlay={overlay_x}:{overlay_y}:format=yuv420[outv]"
        )

        command = [
            self._ffmpeg_bin,
            "-y",
            "-loop",
            "1",
            "-i",
            background_path,
            "-i",
            video_path,
            "-filter_complex",
            filter_complex,
            "-map",
            "[outv]",
            "-map",
            "1:a?",
            "-c:v",
            "libx264",
            "-preset",
            "veryfast",
            "-crf",
            "18",
            "-pix_fmt",
            "yuv420p",
            "-c:a",
            "aac",
            "-b:a",
            "192k",
            "-ar",
            "48000",
            "-movflags",
            "+faststart",
            "-shortest",
            "-r",
            "25",
            segment_path,
        ]
        subprocess.run(command, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        return segment_path, max(duration, 0.0)

    def _measure_video_geometry(self, video_path: str) -> Tuple[int, int]:
        """Return the width and height of a video; (0, 0) when unavailable."""
        cap = cv2.VideoCapture(video_path)
        width = 0
        height = 0
        if cap.isOpened():
            try:
                width = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH) or 0)
                height = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT) or 0)
            finally:
                cap.release()
        else:
            cap.release()
        return max(width, 0), max(height, 0)

    def _measure_video_duration(self, video_path: str) -> float:
        cap = cv2.VideoCapture(video_path)
        duration = 0.0
        if cap.isOpened():
            try:
                fps = cap.get(cv2.CAP_PROP_FPS) or 0.0
                frame_count = cap.get(cv2.CAP_PROP_FRAME_COUNT) or 0.0
                if fps > 0.0 and frame_count > 0.0:
                    duration = float(frame_count / fps)
            finally:
                cap.release()
        else:
            cap.release()
        return max(duration, 0.0)

    def _estimate_slide_timings(self, scripts: Dict[int, str]) -> Dict[int, float]:
        """Estimate slide display time using heuristic speech pacing."""
        base_chars_per_second = 6.21 # ppt切换慢于语速则调高 （6.2ppt有点偏慢，6.3偏快）
        punctuation_pauses = {
            "。": 0.40,
            "？": 0.30,
            "?": 0.30,
            "！": 0.30,
            "!": 0.30,
            "；": 0.30,
            ";": 0.30,
            "，": 0.25,
            ",": 0.25,
            "：": 0.25,
            ":": 0.25,
            "…": 0.35,
        }
        number_bonus = 0.20
        english_bonus = 0.06
        technical_bonus = 0.04
        inhale_time = 0.30
        extra_symbols = {
            "-": 0.20,
            "——": 0.20,
            "~": 0.15,
            "（": 0.15,
            "）": 0.15,
            "(": 0.15,
            ")": 0.15,
            "“": 0.10,
            "”": 0.10,
            "'": 0.10,
            '"': 0.10,
            "<": 0.20,
            ">": 0.20,
        }

        estimations: Dict[int, float] = {}
        for slide_idx, text in scripts.items():
            content = text.strip()
            if not content:
                estimations[slide_idx] = 3.0
                continue

            plain = re.sub(r"\s+", "", content)
            char_count = len(plain)
            base_seconds = char_count / base_chars_per_second if char_count else 0.0

            punctuation_bonus_total = sum(punctuation_pauses.get(ch, 0.0) for ch in content)
            punctuation_bonus_total += sum(extra_symbols.get(ch, 0.0) for ch in content)

            number_hits = re.findall(r"\d+(?:[\.，,]\d+)?", content)
            number_bonus_total = len(number_hits) * number_bonus

            english_tokens = re.findall(r"[A-Za-z]+", content)
            english_bonus_total = 0.0
            for token in english_tokens:
                english_bonus_total += english_bonus
                if token.isupper() or (len(token) > 1 and token[0].isupper()):
                    english_bonus_total += technical_bonus

            sentence_units = [segment for segment in re.split(r"[。！？!?\n]+", content) if segment.strip()]
            inhale_bonus = inhale_time * max(len(sentence_units), 1)

            total = base_seconds + punctuation_bonus_total + number_bonus_total + english_bonus_total + inhale_bonus
            estimations[slide_idx] = max(3.0, total)
        return estimations

    def _normalize_slide_durations(
        self,
        slide_order: List[int],
        raw_timings: Dict[int, float],
        total_duration: float,
    ) -> Dict[int, float]:
        if not slide_order:
            return {}
        raw_values = [max(0.5, raw_timings.get(idx, 3.0)) for idx in slide_order]
        total_raw = sum(raw_values)
        if total_duration <= 0.0:
            fallback = 5.0
            return {idx: fallback for idx in slide_order}
        if total_raw <= 0.0:
            equal_share = max(0.5, total_duration / len(slide_order))
            return {idx: equal_share for idx in slide_order}

        normalized: Dict[int, float] = {idx: base for idx, base in zip(slide_order, raw_values)}
        current_total = sum(normalized.values())
        diff = total_duration - current_total
        if abs(diff) > 0.01:
            last_idx = slide_order[-1]
            original_last = normalized[last_idx]
            adjusted_last = max(0.5, original_last + diff)
            normalized[last_idx] = adjusted_last
            logger.info(
                "[ppt-augment] adjusted last slide duration p%s: %.2fs -> %.2fs (Δ%.2fs)",
                last_idx,
                original_last,
                adjusted_last,
                adjusted_last - original_last,
            )
        return normalized

    def _build_slide_cues(self, slide_order: List[int], durations: Dict[int, float]) -> Dict[int, float]:
        cues: Dict[int, float] = {}
        elapsed = 0.0
        for idx in slide_order:
            cues[idx] = max(0.0, elapsed)
            elapsed += max(0.0, durations.get(idx, 0.0))
        return cues

    def _resolve_overlay_position(
        self,
        slide_dims_in: Tuple[float, float],
        positions: Optional[Dict[int, Dict[str, float]]],
        slide_order: List[int],
    ) -> Dict[str, float]:
        slide_width_in, slide_height_in = slide_dims_in
        # try user-provided position first
        if positions:
            for idx in slide_order:
                pos = positions.get(idx)
                if isinstance(pos, dict):
                    resolved = {
                        "x": float(pos.get("x", 0.0)),
                        "y": float(pos.get("y", 0.0)),
                        "width": float(pos.get("width", 0.0)),
                        "height": float(pos.get("height", 0.0)),
                    }
                    if resolved["width"] > 0 and resolved["height"] > 0:
                        return {
                            "x": max(0.0, min(1.0, resolved["x"])),
                            "y": max(0.0, min(1.0, resolved["y"])),
                            "width": max(0.02, min(1.0, resolved["width"])),
                            "height": max(0.02, min(1.0, resolved["height"])),
                        }

        # fall back to default inches configuration
        left_in, top_in, width_in, height_in = self._video_position
        width_ratio = width_in / slide_width_in if slide_width_in else 0.3
        height_ratio = height_in / slide_height_in if slide_height_in else 0.3
        return {
            "x": left_in / slide_width_in if slide_width_in else 0.55,
            "y": top_in / slide_height_in if slide_height_in else 0.1,
            "width": max(0.02, min(1.0, width_ratio)),
            "height": max(0.02, min(1.0, height_ratio)),
        }

    def _split_combined_video(
        self,
        combined_path: str,
        slide_order: List[int],
        durations: Dict[int, float],
        workdir: str,
    ) -> Dict[int, str]:
        if not self._ffmpeg_bin:
            raise RuntimeError("未检测到 ffmpeg，无法拆分数字人视频。")
        if not os.path.isfile(combined_path):
            raise FileNotFoundError(f"未找到数字人合并视频: {combined_path}")

        segments: Dict[int, str] = {}
        start_time = 0.0
        for slide_idx in slide_order:
            duration = durations.get(slide_idx, 0.0)
            if duration <= 0.1:
                continue
            output_path = os.path.join(workdir, f"slide_{slide_idx:03d}.mp4")
            command = [
                self._ffmpeg_bin,
                "-y",
                "-i",
                combined_path,
                "-ss",
                f"{start_time:.3f}",
                "-t",
                f"{duration:.3f}",
                "-map",
                "0:v:0",
                "-map",
                "0:a?",
                "-c:v",
                "libx264",
                "-preset",
                "veryfast",
                "-crf",
                "18",
                "-pix_fmt",
                "yuv420p",
                "-c:a",
                "aac",
                "-b:a",
                "192k",
                "-ar",
                "48000",
                "-movflags",
                "+faststart",
                output_path,
            ]
            try:
                subprocess.run(command, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
            except subprocess.CalledProcessError as exc:
                logger.warning("[ppt-augment] split combined video failed for slide %s: %s", slide_idx, exc)
                break
            if os.path.isfile(output_path):
                segments[slide_idx] = output_path
            start_time += duration
        return segments

    def _apply_slide_transitions(self, presentation: Presentation, timings: Dict[int, float]) -> None:
        for slide_idx, duration in timings.items():
            if slide_idx < 1 or slide_idx > len(presentation.slides):
                continue
            slide = presentation.slides[slide_idx - 1]
            transition = getattr(slide, "slide_show_transition", None)
            if transition is None:
                logger.warning(
                    "[ppt-augment] 当前 python-pptx 版本不支持 slide_show_transition，跳过自动切换设置 (slide=%s)",
                    slide_idx,
                )
                continue
            transition.advance_on_click = False
            transition.advance_on_time = True
            transition.advance_after_time = max(0.5, float(duration))

    def _export_presentation_video(
        self,
        ppt_path: str,
        slide_durations: Dict[int, float],
        output_path: str,
        fps: int = 25,
        quality: int = 100,
    ) -> Optional[str]:
        if comtypes is None:
            logger.warning("[ppt-augment] PowerPoint COM 不可用，无法导出课程视频。")
            return None
        if not os.path.isfile(ppt_path):
            logger.warning("[ppt-augment] PowerPoint 文件缺失，无法导出课程视频: %s", ppt_path)
            return None

        powerpoint = None
        presentation = None
        initialized = False
        try:
            comtypes.CoInitialize()
            initialized = True
            powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
            powerpoint.Visible = 1
            try:
                powerpoint.WindowState = 2  # ppWindowMinimized
            except Exception:
                pass
            presentation = powerpoint.Presentations.Open(ppt_path, WithWindow=True)

            slide_count = presentation.Slides.Count
            for slide_idx, duration in slide_durations.items():
                if slide_idx < 1 or slide_idx > slide_count:
                    continue
                try:
                    slide = presentation.Slides(slide_idx)
                    transition = slide.SlideShowTransition
                    transition.AdvanceOnClick = False
                    transition.AdvanceOnTime = True
                    transition.AdvanceTime = max(0.5, float(duration))
                except Exception as exc:  # pragma: no cover - COM quirks
                    logger.warning(
                        "[ppt-augment] 设置幻灯片 %s 切换时间失败: %s",
                        slide_idx,
                        exc,
                    )

            try:
                presentation.Save()
            except Exception:
                pass

            try:
                presentation.CreateVideo(
                    output_path,
                    UseTimingsAndNarrations=True,
                    DefaultSlideDuration=5,
                    VertResolution=1080,
                    FramesPerSecond=fps,
                    Quality=quality,
                )
            except Exception as exc:
                logger.warning("[ppt-augment] PowerPoint 导出启动失败: %s", exc)
                return None

            status = getattr(presentation, "CreateVideoStatus", None)
            if status is None:
                logger.warning("[ppt-augment] 当前 PowerPoint 不支持 CreateVideoStatus")
                return None
            deadline = time.time() + 900.0
            while status in (0, 1, 2) and time.time() < deadline:
                time.sleep(0.5)
                status = presentation.CreateVideoStatus
            if status == 3 and os.path.isfile(output_path):
                return output_path
            if status == 4:
                logger.warning("[ppt-augment] PowerPoint 导出失败，状态=4")
                return None
            logger.warning("[ppt-augment] PowerPoint 导出超时或状态异常: %s", status)
            return None
        except Exception as exc:
            logger.warning("[ppt-augment] 导出课程视频失败: %s", exc)
            return None
        finally:
            if presentation is not None:
                try:
                    presentation.Close()
                except Exception:
                    pass
            if powerpoint is not None:
                try:
                    powerpoint.Quit()
                except Exception:
                    pass
            if initialized:
                try:
                    comtypes.CoUninitialize()
                except Exception:
                    pass

    def _overlay_avatar_on_ppt_video(
        self,
        base_video_path: str,
        avatar_video_path: str,
        position: Dict[str, float],
        output_path: str,
    ) -> Optional[str]:
        if not self._ffmpeg_bin:
            raise RuntimeError("未检测到 ffmpeg，无法叠加数字人视频。")
        if not os.path.isfile(base_video_path) or not os.path.isfile(avatar_video_path):
            logger.warning("[ppt-augment] 叠加视频缺失: %s / %s", base_video_path, avatar_video_path)
            return None

        base_width, base_height = self._measure_video_geometry(base_video_path)
        if base_width <= 0 or base_height <= 0:
            logger.warning("[ppt-augment] 无法读取课件视频尺寸，叠加终止。")
            return None

        overlay_w = int(round(max(0.02, min(1.0, position.get("width", 0.3))) * base_width))
        overlay_h = int(round(max(0.02, min(1.0, position.get("height", 0.3))) * base_height))
        overlay_x = int(round(max(0.0, min(1.0, position.get("x", 0.6))) * base_width))
        overlay_y = int(round(max(0.0, min(1.0, position.get("y", 0.1))) * base_height))

        overlay_w = max(2, min(base_width, overlay_w + (overlay_w % 2)))
        overlay_h = max(2, min(base_height, overlay_h + (overlay_h % 2)))
        max_x = max(0, base_width - overlay_w)
        max_y = max(0, base_height - overlay_h)
        overlay_x = max(0, min(max_x, overlay_x))
        overlay_y = max(0, min(max_y, overlay_y))

        filter_complex = (
            f"[0:v]setsar=1[vbg];"
            f"[1:v]scale={overlay_w}:{overlay_h}:force_original_aspect_ratio=increase,"
            f"crop={overlay_w}:{overlay_h},setsar=1[vavatar];"
            f"[vbg][vavatar]overlay={overlay_x}:{overlay_y}:format=yuv420[outv]"
        )

        command = [
            self._ffmpeg_bin,
            "-y",
            "-i",
            base_video_path,
            "-i",
            avatar_video_path,
            "-filter_complex",
            filter_complex,
            "-map",
            "[outv]",
            "-map",
            "1:a?",
            "-c:v",
            "libx264",
            "-preset",
            "veryfast",
            "-crf",
            "18",
            "-pix_fmt",
            "yuv420p",
            "-c:a",
            "aac",
            "-b:a",
            "192k",
            "-ar",
            "48000",
            "-movflags",
            "+faststart",
            "-shortest",
            output_path,
        ]
        try:
            subprocess.run(command, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        except subprocess.CalledProcessError as exc:
            logger.warning("[ppt-augment] 数字人叠加失败: %s", exc)
            return None

        return output_path if os.path.isfile(output_path) else None

    def _build_course_static_segment(
        self,
        background_path: str,
        target_width: int,
        target_height: int,
        workdir: str,
        slide_idx: int,
    ) -> Tuple[str, float]:
        duration = max(0.5, float(self._course_static_duration))
        segment_path = os.path.join(workdir, f"course_seg_{slide_idx:03d}.mp4")
        filter_complex = f"[0:v]scale={target_width}:{target_height},format=yuv420p[outv]"
        command = [
            self._ffmpeg_bin,
            "-y",
            "-loop",
            "1",
            "-i",
            background_path,
            "-f",
            "lavfi",
            "-i",
            "anullsrc=channel_layout=mono:sample_rate=48000",
            "-filter_complex",
            filter_complex,
            "-map",
            "[outv]",
            "-map",
            "1:a",
            "-c:v",
            "libx264",
            "-preset",
            "veryfast",
            "-crf",
            "18",
            "-pix_fmt",
            "yuv420p",
            "-c:a",
            "aac",
            "-b:a",
            "192k",
            "-ar",
            "48000",
            "-movflags",
            "+faststart",
            "-t",
            f"{duration}",
            "-r",
            "25",
            segment_path,
        ]
        subprocess.run(command, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        return segment_path, duration

    def _compose_course_video(
        self,
        ppt_bytes: bytes,
        generated_videos: Dict[int, str],
        applied_positions: Dict[int, Dict[str, float]],
        slide_dims_in: Tuple[float, float],
        total_slides: int,
        workdir: str,
        course_filename: str,
    ) -> Tuple[bytes, Dict[str, object]]:
        if not self._ffmpeg_bin:
            raise RuntimeError("未检测到 ffmpeg，无法生成视频课。")

        try:
            preview_images, _ = generate_slide_previews(ppt_bytes, total_slides)
        except Exception as exc:
            raise RuntimeError(f"导出幻灯片图像失败: {exc}") from exc

        if not preview_images:
            raise RuntimeError("未能生成幻灯片图像，无法生成视频课。")

        background_paths: Dict[int, str] = {}
        for idx, image_base64 in enumerate(preview_images):
            slide_idx = idx + 1
            bg_path = os.path.join(workdir, f"course_slide_{slide_idx:03d}.png")
            with open(bg_path, "wb") as f:
                f.write(base64.b64decode(image_base64))
            background_paths[slide_idx] = bg_path

        slide_width_in, slide_height_in = slide_dims_in
        if slide_width_in <= 0 or slide_height_in <= 0:
            raise RuntimeError("幻灯片尺寸无效，无法生成视频课。")

        target_height = int(round(self._course_base_height))
        if target_height < 2:
            target_height = 1080
        if target_height % 2 != 0:
            target_height += 1
        target_width = int(round(target_height * (slide_width_in / slide_height_in)))
        if target_width < 16:
            target_width = 1280
        if target_width % 2 != 0:
            target_width += 1

        segments: List[str] = []
        total_duration = 0.0
        for slide_idx in range(1, total_slides + 1):
            bg_path = background_paths.get(slide_idx)
            if not bg_path:
                continue
            video_path = generated_videos.get(slide_idx)
            if video_path and slide_idx in applied_positions:
                segment_path, duration = self._build_course_segment_with_avatar(
                    bg_path,
                    video_path,
                    applied_positions[slide_idx],
                    target_width,
                    target_height,
                    workdir,
                    slide_idx,
                )
            else:
                segment_path, duration = self._build_course_static_segment(
                    bg_path,
                    target_width,
                    target_height,
                    workdir,
                    slide_idx,
                )
            segments.append(segment_path)
            total_duration += max(duration, 0.0)

        if not segments:
            raise RuntimeError("未生成任何可用的视频片段，无法合成视频课。")

        concat_path = os.path.join(workdir, "course_concat.txt")
        with open(concat_path, "w", encoding="utf-8") as f:
            for segment in segments:
                safe_segment = segment.replace("\\", "/")
                f.write(f"file '{safe_segment}'\n")

        final_name = os.path.basename(course_filename)
        final_path = os.path.join(workdir, final_name)
        concat_command = [
            self._ffmpeg_bin,
            "-y",
            "-f",
            "concat",
            "-safe",
            "0",
            "-i",
            concat_path,
            "-c",
            "copy",
            "-movflags",
            "+faststart",
            final_path,
        ]
        subprocess.run(concat_command, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)

        with open(final_path, "rb") as f:
            course_bytes = f.read()

        course_meta: Dict[str, object] = {
            "filename": final_name,
            "duration": total_duration,
            "segments": len(segments),
            "resolution": (target_width, target_height),
        }
        return course_bytes, course_meta

    def generate(
        self,
        ppt_bytes: bytes,
        docx_bytes: bytes,
        base_filename: str | None = None,
        positions: Optional[Dict[int, Dict[str, float]]] = None,
        produce_course: bool = False,
        voice_id: Optional[str] = None,
    ) -> Tuple[bytes, Dict[str, object], Optional[bytes]]:
        workdir = tempfile.mkdtemp(prefix="pptaugment_")
        try:
            ppt_path = os.path.join(workdir, "input.pptx")
            with open(ppt_path, "wb") as f:
                f.write(ppt_bytes)

            doc_path = os.path.join(workdir, "script.docx")
            with open(doc_path, "wb") as f:
                f.write(docx_bytes)

            scripts, scripts_meta = self._parse_scripts(doc_path)
            if not scripts:
                raise ValueError("脚本文件中未找到形如 p1/p2 的段落标记。")

            presentation = Presentation(ppt_path)
            total_slides = len(presentation.slides)
            slide_width_in = presentation.slide_width / Inches(1)
            slide_height_in = presentation.slide_height / Inches(1)
            slide_dims_in = (slide_width_in, slide_height_in)

            def build_renderer(session_id: int) -> BaseReal:
                if voice_id:
                    return self._builder(session_id, voice_id)
                return self._builder(session_id)

            slide_order = [idx for idx, _ in sorted(scripts.items()) if 1 <= idx <= total_slides]
            if not slide_order:
                raise RuntimeError("脚本未匹配到任何有效的幻灯片索引。")

            combined_body = "\n\n".join(scripts[idx].strip() for idx in slide_order if scripts[idx].strip())
            ssml_prefix = scripts_meta.get("ssml_prefix") if scripts_meta else None
            ssml_suffix = scripts_meta.get("ssml_suffix") if scripts_meta else None
            combined_text = combined_body
            if combined_body:
                needs_wrapper = bool(ssml_prefix or ssml_suffix)
                if not needs_wrapper:
                    lower_body = combined_body.lower()
                    ssml_markers = (
                        "<break",
                        "<say-as",
                        "<phoneme",
                        "<prosody",
                        "<sub",
                        "<emphasis",
                        "<s>",
                        "<p>",
                        "<audio",
                    )
                    if any(marker in lower_body for marker in ssml_markers):
                        needs_wrapper = True
                if needs_wrapper:
                    parts: List[str] = []
                    parts.append(ssml_prefix or "<speak>")
                    parts.append(combined_body)
                    parts.append(ssml_suffix or "</speak>")
                    combined_text = "\n".join(part for part in parts if part)
            if not combined_text:
                raise RuntimeError("脚本内容为空，无法生成数字人朗读。")

            renderer = DigitalHumanRenderer(self._session_factory, build_renderer)
            combined_video_path = os.path.join(workdir, "combined_avatar.mp4")
            try:
                renderer.speak(combined_text, combined_video_path)
            finally:
                renderer.close()

            if not os.path.isfile(combined_video_path):
                raise RuntimeError("数字人朗读输出缺失，生成失败。")

            combined_duration = self._measure_video_duration(combined_video_path)
            if combined_duration <= 0.0:
                raise RuntimeError("无法解析数字人视频时长，请检查生成的文件。")

            ordered_scripts = {idx: scripts[idx] for idx in slide_order}
            logger.info("[ppt-augment] estimating slide timings for %d slides", len(ordered_scripts))
            raw_timings = self._estimate_slide_timings(ordered_scripts)
            slide_durations = self._normalize_slide_durations(slide_order, raw_timings, combined_duration)
            slide_cues = self._build_slide_cues(slide_order, slide_durations)

            logger.info("[ppt-augment] splitting combined video into %d segments", len(slide_order))
            generated_videos = self._split_combined_video(
                combined_video_path,
                slide_order,
                slide_durations,
                workdir,
            )
            logger.info("[ppt-augment] split finished, generated %d segments", len(generated_videos))
            if len(generated_videos) < len(slide_order):
                missing = sorted(set(slide_order) - set(generated_videos.keys()))
                raise RuntimeError(f"未能拆分全部数字人视频，请检查幻灯片 {missing[:3]} 等的音频时长。")

            overlay_position = self._resolve_overlay_position(slide_dims_in, positions, slide_order)

            timed_ppt_path = os.path.join(workdir, "timed.pptx")
            presentation.save(timed_ppt_path)

            ppt_video_path = os.path.join(workdir, "ppt_base.mp4")
            logger.info("[ppt-augment] exporting PowerPoint video to %s", ppt_video_path)
            exported_video = self._export_presentation_video(timed_ppt_path, slide_durations, ppt_video_path)
            if not exported_video:
                raise RuntimeError("PowerPoint 无法导出课程视频，请确认 PowerPoint 已安装且允许可见窗口运行。")
            logger.info("[ppt-augment] PowerPoint export finished: %s", exported_video)

            final_overlay_path = os.path.join(workdir, "course_overlay.mp4")
            logger.info("[ppt-augment] overlaying avatar video onto PPT base")
            overlay_video = self._overlay_avatar_on_ppt_video(exported_video, combined_video_path, overlay_position, final_overlay_path)
            if not overlay_video:
                raise RuntimeError("叠加数字人视频失败，无法生成课程视频。")
            logger.info("[ppt-augment] overlay finished: %s", overlay_video)

            final_video_duration = self._measure_video_duration(overlay_video) or combined_duration
            final_video_resolution = self._measure_video_geometry(overlay_video)

            presentation = Presentation(timed_ppt_path)
            target_slide_index = slide_order[0]
            embed_slide = presentation.slides[target_slide_index - 1]
            applied_position = self._embed_video(embed_slide, combined_video_path, overlay_position, slide_dims_in)
            applied_positions: Dict[int, Dict[str, float]] = {target_slide_index: applied_position}

            output_path = os.path.join(workdir, "augmented.pptx")
            presentation.save(output_path)

            with open(output_path, "rb") as f:
                result_bytes = f.read()

            safe_name = base_filename or "livetalking-augmented"
            if not safe_name.lower().endswith(".pptx"):
                safe_name += ".pptx"

            final_video_name = os.path.basename(overlay_video)
            meta = {
                "filename": safe_name,
                "total_slides": total_slides,
                "video_slides": [target_slide_index],
                "slide_durations": slide_durations,
                "slide_cues": slide_cues,
                "combined_duration": combined_duration,
                "final_video": final_video_name,
                "final_video_duration": final_video_duration,
                "final_video_resolution": final_video_resolution,
                "overlay_position": overlay_position,
            }
            session_id = self._persist_session_assets(
                ppt_bytes,
                generated_videos,
                applied_positions,
                slide_dims_in,
                total_slides,
                voice_id,
                slide_durations,
                slide_cues,
                combined_video_path,
                overlay_video,
                {
                    "duration": final_video_duration,
                    "resolution": final_video_resolution,
                    "position": overlay_position,
                    "slide": target_slide_index,
                },
            )
            meta["session_id"] = session_id
            course_bytes: Optional[bytes] = None
            meta.update(
                {
                    "course_filename": final_video_name,
                    "course_duration": final_video_duration,
                    "course_segments": 1,
                    "course_resolution": final_video_resolution,
                }
            )
            if produce_course:
                with open(overlay_video, "rb") as course_file:
                    course_bytes = course_file.read()

            meta["voice"] = voice_id
            return result_bytes, meta, course_bytes
        finally:
            shutil.rmtree(workdir, ignore_errors=True)
